#!/usr/bin/env python3
"""
Create PowerPoint presentations from outlines, topics, or data using templates.
Enterprise-grade features: template management, intelligent layouts, quality checks.
"""

import argparse
import json
import sys
import logging
from pathlib import Path
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from template_manager import get_template, get_template_pptx_path, list_templates
from config import get_config
from content_enhancer import ContentEnhancer

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PresentationCreator:
    def __init__(self, config=None):
        self.config = config or get_config()
        self.enhancer = ContentEnhancer()
        self.current_prs = None
        self.template_info = None
    
    def create_with_template(self, prs_data: Dict[str, Any], template_id: str, output_path: str) -> int:
        self.template_info = get_template(template_id)
        if not self.template_info:
            logger.warning(f"Template not found: {template_id}")
            return self.create_default(prs_data, output_path)
        
        template_pptx_path = get_template_pptx_path(template_id)
        if not template_pptx_path or not Path(template_pptx_path).exists():
            return self.create_default(prs_data, output_path)
        
        logger.info(f"Using template: {self.template_info['name']}")
        self.current_prs = Presentation(template_pptx_path)
        
        for i in range(len(self.current_prs.slides) - 1, -1, -1):
            rId = self.current_prs.slides._sldIdLst[i].rId
            self.current_prs.part.drop_rel(rId)
            del self.current_prs.slides._sldIdLst[i]
        
        self._build_slides(prs_data)
        self.current_prs.save(output_path)
        return len(self.current_prs.slides)
    
    def create_default(self, prs_data: Dict[str, Any], output_path: str) -> int:
        self.current_prs = Presentation()
        self._build_slides(prs_data)
        self.current_prs.save(output_path)
        return len(self.current_prs.slides)
    
    def _build_slides(self, prs_data: Dict[str, Any]):
        title = prs_data.get("title", "Presentation")
        subtitle = prs_data.get("subtitle", "")
        self._add_title_slide(title, subtitle)
        for slide_data in prs_data.get("slides", []):
            self._add_content_slide(slide_data)
    
    def _add_title_slide(self, title: str, subtitle: str = ""):
        if self.current_prs.slide_layouts:
            try:
                slide = self.current_prs.slides.add_slide(self.current_prs.slide_layouts[0])
                if slide.shapes.title:
                    slide.shapes.title.text = title
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle
                return
            except: pass
        slide = self.current_prs.slides.add_slide(self.current_prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
    
    def _add_content_slide(self, slide_data: Dict[str, Any]):
        layout_type = slide_data.get("layout", "content")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        
        if layout_type == "title":
            self._add_title_slide(title, "")
        elif layout_type == "section":
            self._add_section_slide(title)
        else:
            self._add_standard_slide(title, bullets)
    
    def _add_standard_slide(self, title: str, bullets: list):
        if len(self.current_prs.slide_layouts) > 1:
            try:
                slide = self.current_prs.slides.add_slide(self.current_prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = title
                if bullets:
                    for placeholder in slide.placeholders:
                        if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type == 7:
                            tf = placeholder.text_frame
                            tf.clear()
                            for i, bullet in enumerate(bullets):
                                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                                p.text = bullet
                            break
                return
            except: pass
        slide = self.current_prs.slides.add_slide(self.current_prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = title
    
    def _add_section_slide(self, title: str):
        slide = self.current_prs.slides.add_slide(self.current_prs.slide_layouts[6])
        left = top = Inches(0)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, self.current_prs.slide_width, self.current_prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

def parse_markdown_outline(content: str) -> Dict[str, Any]:
    lines = content.strip().split('\n')
    presentation = {"title": "", "subtitle": "", "slides": []}
    current_slide = None
    current_bullets = []
    for line in lines:
        line = line.strip()
        if line.startswith('subtitle:'):
            presentation["subtitle"] = line.split(':', 1)[1].strip()
        elif line.startswith('# ') and not presentation["title"]:
            presentation["title"] = line[2:].strip()
        elif line.startswith('## '):
            if current_slide:
                current_slide["bullets"] = current_bullets
                presentation["slides"].append(current_slide)
            current_slide = {"title": line[3:].strip(), "bullets": []}
            current_bullets = []
        elif line.startswith('- '):
            current_bullets.append(line[2:].strip())
    if current_slide:
        current_slide["bullets"] = current_bullets
        presentation["slides"].append(current_slide)
    return presentation

def generate_topic_slides(topic: str, num_slides: int = 5) -> list:
    return [
        {"title": topic, "layout": "title", "bullets": []},
        {"title": "Overview", "layout": "content", "bullets": [f"Introduction to {topic}", "Key objectives", "What we'll cover"]},
        {"title": "Key Points", "layout": "content", "bullets": ["Main concept #1", "Main concept #2", "Main concept #3"]},
        {"title": "Analysis", "layout": "content", "bullets": ["Data and insights", "Trends"]},
        {"title": "Conclusion", "layout": "content", "bullets": ["Summary", "Next steps", "Q&A"]}
    ][:num_slides]

def main():
    parser = argparse.ArgumentParser(description='Create PowerPoint presentations')
    parser.add_argument("--outline", help="Markdown outline file")
    parser.add_argument("--topic", help="Topic for auto-generation")
    parser.add_argument("--slides", type=int, default=5)
    parser.add_argument("--json", help="JSON structure file")
    parser.add_argument("--template", "-t", help="Template ID/name")
    parser.add_argument("--list-templates", action="store_true")
    parser.add_argument("--output", "-o", required=True)
    parser.add_argument("--verbose", "-v", action="store_true")
    parser.add_argument("--quality-check", "-q", action="store_true")
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    if args.list_templates:
        templates = list_templates()
        if not templates:
            print("📭 No templates found")
            return 0
        for t in templates:
            print(f"📽️ {t['name']} (ID: {t['id']}) - {t['slide_count']} slides")
        return 0
    
    prs_data = {"title": "Presentation", "subtitle": "", "slides": []}
    if args.outline:
        with open(args.outline, 'r', encoding='utf-8') as f:
            prs_data = parse_markdown_outline(f.read())
    elif args.topic:
        prs_data["title"] = args.topic
        prs_data["slides"] = generate_topic_slides(args.topic, args.slides)
    elif args.json:
        with open(args.json, 'r', encoding='utf-8') as f:
            prs_data = json.load(f)
    else:
        print("Error: Specify --outline, --topic, or --json")
        return 1
    
    if args.quality_check:
        enhancer = ContentEnhancer()
        print("\n📊 Quality Check:")
        for slide in prs_data.get('slides', []):
            q = enhancer.check_content_quality(slide.get('title', ''), slide.get('bullets', []))
            if q['score'] < 100:
                print(f"  ⚠️ {slide.get('title')}: {q['score']}/100 - {q['suggestions']}")
    
    creator = PresentationCreator()
    try:
        count = creator.create_with_template(prs_data, args.template, args.output) if args.template else creator.create_default(prs_data, args.output)
        print(f"\n✅ Created {args.output} ({count} slides)")
        return 0
    except Exception as e:
        logger.error(f"Failed: {e}")
        print(f"❌ Error: {e}")
        return 1

if __name__ == "__main__":
    sys.exit(main())
