#!/usr/bin/env python3
"""
Analyze a PowerPoint file and extract template information for reuse.
Extracts colors, fonts, layouts, logos, and design elements.
"""

import argparse
import json
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Inches

def extract_colors_from_theme(theme_xml):
    """Extract color scheme from theme XML."""
    colors = {}
    try:
        for elem in theme_xml.iter():
            if 'clrMap' in elem.tag:
                attrs = elem.attrib
                colors['bg1'] = attrs.get('bg1', 'FFFFFF')
                colors['tx1'] = attrs.get('tx1', '000000')
                colors['bg2'] = attrs.get('bg2', 'FFFFFF')
                colors['tx2'] = attrs.get('tx2', '000000')
                colors['accent1'] = attrs.get('accent1', '0000FF')
                colors['accent2'] = attrs.get('accent2', 'FF00FF')
                colors['accent3'] = attrs.get('accent3', '00FFFF')
                colors['accent4'] = attrs.get('accent4', '800080')
                colors['accent5'] = attrs.get('accent5', '008000')
                colors['accent6'] = attrs.get('accent6', 'FF0000')
    except Exception as e:
        colors['error'] = str(e)
    
    return colors

def extract_fonts_from_theme(theme_xml):
    """Extract font scheme from theme XML."""
    fonts = {'heading': 'Arial', 'body': 'Arial'}
    try:
        for elem in theme_xml.iter():
            if 'fontScheme' in elem.tag or 'fontCollection' in elem.tag:
                for child in elem:
                    if 'majorFont' in child.tag:
                        for font in child:
                            if 'latin' in font.tag:
                                fonts['heading'] = font.attrib.get('typeface', 'Arial')
                    elif 'minorFont' in child.tag:
                        for font in child:
                            if 'latin' in font.tag:
                                fonts['body'] = font.attrib.get('typeface', 'Arial')
    except Exception:
        pass
    
    return fonts

def analyze_slide_layouts(prs):
    """Analyze slide layouts in the presentation."""
    layouts = []
    
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            'index': i,
            'name': layout.name if hasattr(layout, 'name') else f'Layout {i}',
        }
        
        # Analyze placeholders
        placeholders = []
        for placeholder in layout.placeholders:
            ph_info = {
                'id': placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else i,
                'type': str(placeholder.placeholder_format.type) if hasattr(placeholder, 'placeholder_format') else 'Unknown',
                'name': placeholder.name if hasattr(placeholder, 'name') else 'Unknown',
            }
            placeholders.append(ph_info)
        
        layout_info['placeholders'] = placeholders
        layouts.append(layout_info)
    
    return layouts

def analyze_master_slides(prs):
    """Analyze master slides for background and design elements."""
    masters = []
    
    for i, master in enumerate(prs.slide_masters):
        master_info = {
            'index': i,
            'name': master.name if hasattr(master, 'name') else f'Master {i}',
        }
        
        # Background
        if master.background:
            try:
                fill = master.background.fill
                if fill.solid:
                    if fill.fore_color.rgb:
                        master_info['background_color'] = str(fill.fore_color.rgb)
                    elif fill.fore_color.theme_color:
                        master_info['background_theme_color'] = str(fill.fore_color.theme_color)
            except Exception:
                pass
        
        # Logos and images on master
        images = []
        for shape in master.shapes:
            if shape.shape_type == 13:  # Picture
                images.append({
                    'type': 'logo' if 'logo' in shape.name.lower() else 'image',
                    'name': shape.name,
                    'position': {
                        'x': str(shape.left),
                        'y': str(shape.top),
                        'width': str(shape.width),
                        'height': str(shape.height)
                    }
                })
        
        master_info['master_images'] = images
        masters.append(master_info)
    
    return masters

def extract_sample_colors(prs):
    """Extract commonly used colors from the presentation."""
    color_counts = {}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'fill') and shape.fill:
                try:
                    if shape.fill.solid and shape.fill.fore_color:
                        if shape.fill.fore_color.rgb:
                            color = str(shape.fill.fore_color.rgb)
                            color_counts[color] = color_counts.get(color, 0) + 1
                except Exception:
                    pass
            
            # Text colors
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                color = str(run.font.color.rgb)
                                color_counts[color] = color_counts.get(color, 0) + 1
                        except Exception:
                            pass
    
    # Sort by frequency
    sorted_colors = sorted(color_counts.items(), key=lambda x: x[1], reverse=True)
    
    return {
        'most_used': [c[0] for c in sorted_colors[:5]],
        'all_colors': list(color_counts.keys())
    }

def analyze_template(pptx_path, output_json=None):
    """Main analysis function."""
    pptx_path = Path(pptx_path)
    
    if not pptx_path.exists():
        return {'error': f'File not found: {pptx_path}'}
    
    try:
        prs = Presentation(pptx_path)
        
        template_info = {
            'source_file': str(pptx_path.name),
            'slide_count': len(prs.slides),
            'layout_count': len(prs.slide_layouts),
            'master_count': len(prs.slide_masters),
        }
        
        # Analyze layouts
        template_info['layouts'] = analyze_slide_layouts(prs)
        
        # Analyze masters
        template_info['masters'] = analyze_master_slides(prs)
        
        # Extract colors
        template_info['sample_colors'] = extract_sample_colors(prs)
        
        # Try to extract theme
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                if 'ppt/theme/theme1.xml' in zf.namelist():
                    theme_xml = ET.parse(zf.open('ppt/theme/theme1.xml'))
                    template_info['theme_colors'] = extract_colors_from_theme(theme_xml)
                    template_info['theme_fonts'] = extract_fonts_from_theme(theme_xml)
        except Exception as e:
            template_info['theme_error'] = str(e)
        
        # Save to JSON if requested
        if output_json:
            with open(output_json, 'w', encoding='utf-8') as f:
                json.dump(template_info, f, indent=2, default=str)
        
        return template_info
        
    except Exception as e:
        return {'error': str(e)}

def main():
    parser = argparse.ArgumentParser(description='Analyze PowerPoint template')
    parser.add_argument('input', help='Input PPTX file')
    parser.add_argument('--output', '-o', help='Output JSON file')
    parser.add_argument('--json', action='store_true', help='Output as JSON (default: human-readable)')
    
    args = parser.parse_args()
    
    result = analyze_template(args.input, args.output)
    
    if args.json or not args.output:
        print(json.dumps(result, indent=2, default=str))
    else:
        # Human-readable output
        print(f"\n📊 Template Analysis: {result.get('source_file', 'Unknown')}")
        print(f"   Slides: {result.get('slide_count', 0)}")
        print(f"   Layouts: {result.get('layout_count', 0)}")
        print(f"   Masters: {result.get('master_count', 0)}")
        
        if 'theme_fonts' in result and result['theme_fonts']:
            print(f"\n📝 Fonts:")
            print(f"   Heading: {result['theme_fonts'].get('heading', 'N/A')}")
            print(f"   Body: {result['theme_fonts'].get('body', 'N/A')}")
        
        if 'sample_colors' in result and result['sample_colors'].get('most_used'):
            print(f"\n🎨 Top Colors: {', '.join(result['sample_colors'].get('most_used', [])[:5])}")
        
        if args.output:
            print(f"\n✅ Template saved to: {args.output}")
    
    return 0 if 'error' not in result else 1

if __name__ == '__main__':
    sys.exit(main())
